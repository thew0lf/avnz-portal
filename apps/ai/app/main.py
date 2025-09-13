import hashlib, io, os, uuid, re, json, math, hmac, base64
from typing import List, Optional, Tuple
from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from pydantic import BaseModel
import psycopg2, psycopg2.extras
import tiktoken
from pdfminer.high_level import extract_text as pdf_extract
from docx import Document as Docx
from pptx import Presentation
from openpyxl import load_workbook

app = FastAPI(title="avnzr-ai")

DB_URL = os.environ.get("DATABASE_URL")
PGSSL = os.environ.get("PGSSL")
USE_BEDROCK = os.environ.get("USE_BEDROCK","false").lower()=="true"
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
OPENAI_PROJECT = os.environ.get("OPENAI_PROJECT")
AWS_REGION = os.environ.get("AWS_REGION","us-east-1")
AUTH_SECRET = os.environ.get("AUTH_SECRET","dev-secret-change-me")

def pg():
    sslmode = 'require' if PGSSL else None
    conn = psycopg2.connect(DB_URL, sslmode=sslmode) if sslmode else psycopg2.connect(DB_URL)
    return conn

def is_project_member(project_id: Optional[str], user_id: Optional[str]) -> bool:
    if not project_id or not user_id:
        return False
    conn = pg()
    try:
        with conn.cursor() as cur:
            cur.execute("select 1 from project_members where project_id=%s and user_id=%s limit 1", (project_id, user_id))
            return cur.fetchone() is not None
    except Exception:
        return False
    finally:
        conn.close()

def pii_counts(text: str) -> Tuple[int,int,int,int]:
    email = len(re.findall(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", text))
    phone = len(re.findall(r"\b(?:\+?1[-.\s]?)?(?:\(?\d{3}\)?[-.\s]?)?\d{3}[-.\s]?\d{4}\b", text))
    ssn = len(re.findall(r"\b\d{3}-\d{2}-\d{4}\b", text))
    cc = len(re.findall(r"\b(?:\d[ -]*?){13,16}\b", text))
    return email, phone, ssn, cc

def redact(text: str) -> str:
    text = re.sub(r"([A-Za-z0-9._%+-]+)@([A-Za-z0-9.-]+\.[A-Za-z]{2,})", r"***@\2", text)
    text = re.sub(r"\b(\+?1[-.\s]?)?(\(?\d{3}\)?[-.\s]?)?\d{3}[-.\s]?\d{4}\b", "[phone]", text)
    text = re.sub(r"\b\d{3}-\d{2}-\d{4}\b", "[ssn]", text)
    text = re.sub(r"\b(?:\d[ -]*?){13,16}\b", "[cc]", text)
    return text

def tokenc(text: str) -> int:
    enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))

def hash_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()
def b64url_decode(s: str) -> bytes:
    s = s.replace('-', '+').replace('_', '/')
    pad = len(s) % 4
    if pad:
        s += '=' * (4 - pad)
    return base64.b64decode(s)
def verify_token(token: str, secret: str):
    try:
        p, sig = token.split('.')
        expected = hmac.new(secret.encode('utf-8'), p.encode('utf-8'), hashlib.sha256).digest()
        # handle missing padding in base64url
        pad = '=' * (-len(sig) % 4)
        got = base64.urlsafe_b64decode(sig + pad)
        if not hmac.compare_digest(got, expected):
            return None
        payload = json.loads(b64url_decode(p).decode('utf-8'))
        return payload
    except Exception:
        return None

def openai_embed(texts: List[str]) -> List[List[float]]:
    from openai import OpenAI
    # Prefer project scoping if available; fall back safely if client doesn't support it
    try:
        client = OpenAI(api_key=OPENAI_API_KEY, project=OPENAI_PROJECT) if OPENAI_PROJECT else OpenAI(api_key=OPENAI_API_KEY)
    except TypeError:
        client = OpenAI(api_key=OPENAI_API_KEY)
    resp = client.embeddings.create(model="text-embedding-3-small", input=texts)
    return [d.embedding for d in resp.data]

def bedrock_embed(texts: List[str]) -> List[List[float]]:
    import boto3, json
    rt = boto3.client("bedrock-runtime", region_name=AWS_REGION)
    out = []
    for t in texts:
        body = json.dumps({"inputText": t})
        r = rt.invoke_model(modelId="amazon.titan-embed-text-v2:0", body=body)
        payload = json.loads(r.get("body").read())
        out.append(payload.get("embedding", []))
    return out

def embed(texts: List[str]) -> List[List[float]]:
    if USE_BEDROCK:
        return bedrock_embed(texts)
    return openai_embed(texts)

def chunk_text(text: str, target_tokens: int = 500, overlap: int = 50) -> List[str]:
    enc = tiktoken.get_encoding("cl100k_base")
    ids = enc.encode(text)
    chunks = []
    i = 0
    while i < len(ids):
        j = min(i + target_tokens, len(ids))
        chunk = enc.decode(ids[i:j])
        chunks.append(chunk)
        i = j - overlap
        if i < 0: i = 0
    return [c for c in chunks if c.strip()]

class SearchReq(BaseModel):
    org_id: str
    query: str
    k: int = 5
    project_id: Optional[str] = None
    project_code: Optional[str] = None

@app.get("/healthz")
def healthz():
    return {"ok": True, "provider": ("bedrock" if USE_BEDROCK else "openai")}

@app.post("/ingest")
async def ingest(org_id: Optional[str] = Form(None), user_id: Optional[str] = Form(None), plan: str = Form("pro"), project_code: Optional[str] = Form(None), project_id: Optional[str] = Form(None), file: UploadFile = File(...), authorization: Optional[str] = None):
    # Auth
    authz = authorization or ""
    token = authz[7:] if authz.lower().startswith("bearer ") else None
    sess = verify_token(token or "", AUTH_SECRET)
    if not sess:
        raise HTTPException(status_code=401, detail="unauthorized")
    perms = sess.get("perms") or []
    if "ingest" not in perms and "admin" not in perms:
        raise HTTPException(status_code=403, detail="forbidden")
    org_id = sess.get("orgUUID") or sess.get("orgId")
    user_id = sess.get("userId")
    data = await file.read()
    text = ""
    mt = file.content_type or "application/octet-stream"
    try:
        if mt == "application/pdf" or file.filename.lower().endswith(".pdf"):
            text = pdf_extract(io.BytesIO(data))
        elif mt in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or file.filename.lower().endswith(".docx"):
            doc = Docx(io.BytesIO(data))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif mt in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",) or file.filename.lower().endswith(".pptx"):
            pres = Presentation(io.BytesIO(data))
            text = "\n".join(shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif mt in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",) or file.filename.lower().endswith(".xlsx"):
            wb = load_workbook(io.BytesIO(data))
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join([str(x) for x in row if x is not None]))
            text = "\n".join(parts)
        elif mt.startswith("text/") or file.filename.lower().endswith(('.txt','.md','.csv','.json')):
            text = data.decode("utf-8", errors="ignore")
        elif mt.startswith("image/"):
            # OCR not enabled by default; store placeholder
            text = f"[image:{file.filename}]"
        else:
            text = data.decode("utf-8", errors="ignore")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"failed to parse: {e}")

    # resolve project id if code provided
    pid = None
    if project_id:
        pid = project_id
    elif project_code:
        try:
            conn0 = pg()
            with conn0.cursor() as cur:
                cur.execute("select id from projects where org_id=%s and code=%s", (org_id, project_code))
                row = cur.fetchone(); pid = row[0] if row else None
        except Exception:
            pid = None
        finally:
            try:
                conn0.close()
            except Exception:
                pass

    # enforce project membership if pid specified
    perms = sess.get("perms") or []
    if pid and ("admin" not in perms and "manage_projects" not in perms):
        if not is_project_member(pid, user_id):
            raise HTTPException(status_code=403, detail="forbidden: project membership required")

    doc_id = str(uuid.uuid4())
    conn = pg(); conn.autocommit = False
    try:
        with conn.cursor() as cur:
            cur.execute("insert into documents(id, org_id, project_id, user_id, filename, mime_type, size_bytes, plan) values (%s,%s,%s,%s,%s,%s,%s,%s)", (doc_id, org_id, pid, user_id, file.filename, mt, len(data), plan))
        chunks = chunk_text(text, 600, 80)
        red_counts_total = [0,0,0,0]
        redacted_chunks = []
        for idx, ck in enumerate(chunks):
            rc = pii_counts(ck)
            red_counts_total = [a+b for a,b in zip(red_counts_total, rc)]
            red = redact(ck)
            redacted_chunks.append((idx, red, rc))
        # embeddings (with cache)
        to_embed = [red for _, red, _ in redacted_chunks]
        hashes = [hash_text(t) for t in to_embed]
        with conn.cursor(cursor_factory=psycopg2.extras.DictCursor) as cur:
            cur.execute("select content_hash, embedding from embeddings_cache where content_hash = any(%s)", (hashes,))
            cached = {r[0]: r[1] for r in cur.fetchall()}
        need = []
        order = []
        for h, t in zip(hashes, to_embed):
            if h in cached: order.append(cached[h])
            else: need.append(t); order.append(None)
        new_embs = embed(need) if need else []
        iter_new = iter(new_embs)
        embeddings = []
        with conn.cursor() as cur:
            for idx, (h, t) in enumerate(zip(hashes, to_embed)):
                if order[idx] is None:
                    emb = next(iter_new)
                    embeddings.append(emb)
                    cur.execute("insert into embeddings_cache(content_hash, embedding, provider, model) values (%s,%s,%s,%s) on conflict (content_hash) do nothing",
                        (h, emb, "bedrock" if USE_BEDROCK else "openai", "titan-v2" if USE_BEDROCK else "text-embedding-3-small"))
                else:
                    embeddings.append(order[idx])
        with conn.cursor() as cur:
            for (idx, red, rc), emb in zip(redacted_chunks, embeddings):
                cur.execute("insert into document_chunks(id, document_id, org_id, project_id, chunk_index, content, content_redacted, embedding) values (gen_random_uuid(), %s, %s, %s, %s, %s, %s, %s)",
                            (doc_id, org_id, pid, idx, chunks[idx], red, emb))
        # redaction event (aggregate)
        with conn.cursor() as cur:
            cur.execute("insert into redaction_events(node_type, node_id, document_id, pii_email, pii_phone, pii_ssn, pii_cc) values ('document', %s, %s, %s, %s, %s, %s)",
                        (doc_id, doc_id, red_counts_total[0], red_counts_total[1], red_counts_total[2], red_counts_total[3]))
        conn.commit()
        return {"document_id": doc_id, "chunks": len(chunks), "pii": {"email": red_counts_total[0], "phone": red_counts_total[1], "ssn": red_counts_total[2], "cc": red_counts_total[3]}}
    except Exception as e:
        conn.rollback(); raise
    finally:
        conn.close()

@app.post("/search")
def search(req: SearchReq, authorization: Optional[str] = None):
    authz = authorization or ""
    token = authz[7:] if authz.lower().startswith("bearer ") else None
    sess = verify_token(token or "", AUTH_SECRET)
    if not sess:
        raise HTTPException(status_code=401, detail="unauthorized")
    perms = sess.get("perms") or []
    if "search" not in perms and "admin" not in perms:
        raise HTTPException(status_code=403, detail="forbidden")
    # Enforce org from token (prefer UUID)
    req.org_id = sess.get("orgUUID") or sess.get("orgId")
    # resolve project
    pid = req.project_id
    if (not pid) and req.project_code:
        try:
            conn0 = pg()
            with conn0.cursor() as cur:
                cur.execute("select id from projects where org_id=%s and code=%s", (req.org_id, req.project_code))
                row = cur.fetchone(); pid = row[0] if row else None
        except Exception:
            pid = None
        finally:
            try:
                conn0.close()
            except Exception:
                pass

    # enforce project membership if pid specified
    perms = sess.get("perms") or []
    if pid and ("admin" not in perms and "manage_projects" not in perms):
        if not is_project_member(pid, sess.get("userId")):
            raise HTTPException(status_code=403, detail="forbidden: project membership required")

    qemb = embed([req.query])[0]
    conn = pg()
    try:
        with conn.cursor(cursor_factory=psycopg2.extras.DictCursor) as cur:
            if pid:
                cur.execute("""
                    select document_id, chunk_index, content_redacted, 1.0/(1.0 + (embedding <-> %s)) as score
                    from document_chunks where org_id = %s and project_id = %s
                    order by embedding <-> %s asc
                    limit %s
                """ , (qemb, req.org_id, pid, qemb, req.k))
            else:
                cur.execute("""
                    select document_id, chunk_index, content_redacted, 1.0/(1.0 + (embedding <-> %s)) as score
                    from document_chunks where org_id = %s
                    order by embedding <-> %s asc
                    limit %s
                """ , (qemb, req.org_id, qemb, req.k))
            rows = [dict(r) for r in cur.fetchall()]
            return {"matches": rows}
    finally:
        conn.close()
